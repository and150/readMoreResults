import constants as cts
import statistics
from datetime import datetime, timedelta
from pathlib import Path

from mpl_toolkits.axes_grid1 import host_subplot
import mpl_toolkits.axisartist as AA
import matplotlib.dates as md
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.legend_handler import HandlerPatch

import io, re #,os
from openpyxl import load_workbook, Workbook

from pandas.plotting import register_matplotlib_converters
register_matplotlib_converters()


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

### local constants ###
#STAT_FILE = "D:/home/py/excel_parse/model_stat.xlsx"
#STAT_FILTER = '.+' # all wells
STAT_FILE = "D:/WQ2/HM/hm_journal.xlsx"
STAT_FILTER = '(?!WQ2-137)(WQ2-.+)|(WQ-11)|(WQ-13)' # all WQ2- and two WQ wells

STAT_WORKSHEET = 'stat_ext'
PRS_TEMPLATE = "D:/WQ2/HM/prs_template.pptx"
TEMPLATE_NUMBER  = 4  # the number of layout in user templates (check in Slide Master)
TB_WIDTH = 23
TB_START_ROW = 23


def get_statistics(statistics_file_name, well_names):
    wb = load_workbook(statistics_file_name, data_only=True)
    ws = wb[STAT_WORKSHEET]  
    header = [ws.cell(TB_START_ROW,col).value for col in range(1,TB_WIDTH)] # no replacing Delta
    #header = [ ''.join( 'd' if x=='\u0394' else x for x in list(ws.cell(TB_START_ROW,col).value)) for col in range(1,TB_WIDTH)]

    well_matched = True
    data_table = {} 
    ind = {4:5, 8:10, 11:10, 14:10, 17:10, 20:10, 21:10} # {column_indexes:error_tolerances} for hm_journal.xlsx !!!
    row = TB_START_ROW+1
    while ws.cell(row, 1).value:
        error_colours = ['none']*(TB_WIDTH-1)
        data_record =[0]*(TB_WIDTH-1)

        if ws.cell(row,1).value in well_names:
            data_record = [ws.cell(row,col).value for col in range(1,TB_WIDTH)]

            for i in ind:
                if isinstance(data_record[i], float) or isinstance(data_record[i], int):
                    if data_record[i] > ind[i]: 
                        error_colours[i] = 'lightcoral'
                        well_matched = False
                    elif data_record[i] < ind[i]*(-1):
                        error_colours[i] = 'lightblue'
                        well_matched = False
                    else:
                        error_colours[i] = 'none'
            data_table.update({data_record[0]:[[f"{x:.1f}" if isinstance(x,float) or isinstance(x,int) else x for x in data_record], error_colours, well_matched]})
            well_matched = True
        row+=1
    return (header, data_table)



def make_graph(root_name, well_name, x_values, ys_values, header, well_stat, prs):
    fig = plt.figure()
    fig.patch.set_facecolor('b')
    fig.patch.set_alpha(0.5)
    #print(well_stat)
    host = host_subplot(111, axes_class=AA.Axes)
    #plt.subplots(figsize=(20,10))
    plt.subplots_adjust(right=1.0)

    par1 = host.twinx()
    par2 = host.twinx()
    par3 = host.twinx()

    new_fixed_axis = par2.get_grid_helper().new_fixed_axis
    par2.axis["right"] = new_fixed_axis(loc="right",axes=par2, offset=(50,0))

    new_fixed_axis = par3.get_grid_helper().new_fixed_axis
    par3.axis["right"] = new_fixed_axis(loc="right",axes=par3, offset=(100,0))

    par1.axis["right"].toggle(all=True)
    par2.axis["right"].toggle(all=True)
    par3.axis["right"].toggle(all=True)

    host.get_xaxis().set_major_locator(md.YearLocator())
    host.get_xaxis().set_major_formatter(md.DateFormatter("%d.%m.%Y"))

    # TODO make function to calculate limits (check defaults, than calculate custom)
    pres_lim = max(max(ys_values['Sbhp']), max(ys_values['Hbhp']))
    if pres_lim > 500: pres_lim = 500 
    liq_lim = max(max(ys_values['Sliq']), max(ys_values['Hliq']), max(ys_values['Swir']), max(ys_values['Hwir']))
    wPI_lim = max(ys_values['wPI4']) 
    if wPI_lim > 10000: wPI_lim = statistics.median(ys_values['wPI4']) + 1000 

    host.set_xlim(x_values[next((x for x,y in enumerate(ys_values['Hliq']) if y) ,0)], x_values[-1]) # set datetime limits
    host.set_ylim(0,pres_lim) # pressure
    par1.set_ylim(0,liq_lim) # liquid and injection
    par2.set_ylim(0,100) # wcut
    par3.set_ylim(0,wPI_lim) # wPI

    #plt.title(f"{well_name}")
    host.set_xlabel("Date")
    host.set_ylabel("P, barsa")
    par1.set_ylabel("Q, sm3/day")
    par2.set_ylabel("WCT, %")
    par3.set_ylabel("WPI")

    host.plot(x_values, ys_values['Sbhp'], '-', color='black', label="WBHP", linewidth=0.8, zorder=15) # BHP
    host.plot(x_values, ys_values['Hbhp'], '.', markersize=4, color='red', label="WBHPH", zorder=5) # BHPH

    par1.plot(x_values, ys_values['Sliq'], '-', color='lime', label="WLPR", linewidth=0.8, zorder=10) # LIQ
    par1.plot(x_values, ys_values['Hliq'], '.', markersize=3, color='green', label="WLPRH", zorder=0) # LIQH

    par2.plot(x_values, ys_values['Swcut'], '-', color='cyan', label="WWCT", linewidth=0.8, zorder=10) # WCUT
    par2.plot(x_values, ys_values['Hwcut'], '.', markersize=3, color='cyan', label="WWCTH", zorder=0) # WCUTH

    par1.plot_date(x_values, ys_values['Swir'], '-', color='blue', label="WWIN", linewidth=0.8, zorder=10) # WINJ
    par1.plot_date(x_values, ys_values['Hwir'], '.', markersize=2, marker='s', color='blue', label="WWINH", zorder=0) # WINJH

    par3.plot(x_values, ys_values['wPI4'], '-', color='magenta', label="wPI4", linewidth=0.5, zorder=5) # wPI4

    host.legend(bbox_to_anchor=(0.0, -0.4, 1.0, 0.3), ncol=5, mode="expand", columnspacing=1.0) # dimensions hardcoded :(

    # TODO think how better pass the statistics file name into the program, now it is a constant!
    # TODO find out how to convert all interface constants into proportions
    #!!!constant columns widths for 22 columns hm_journal.xlsx format !!!
    col_width = [0.047,0.056,0.056,0.044,0.044,0.050,0.050,0.044,0.044,0.050,0.050,0.044,0.039,0.039,0.044,0.051,0.051,0.038,0.039,0.039,0.039,0.040]
    the_table = plt.table(cellText=[well_stat[0]], 
                colLabels=header, 
                colWidths=col_width,
                colColours=['lightgrey']*(TB_WIDTH-1),
                cellColours=[well_stat[1]],
                cellLoc='center',
                bbox=(-0.10, -0.50, 2.00, 0.23)) # dimensions hardcoded :(
    the_table.auto_set_font_size(False)
    the_table.set_fontsize(7.6)
    for col in range(TB_WIDTH-1):
        the_table[0,col].set_height(0.8)
        the_table[1,col].set_height(0.2)

    #plt.savefig('./'+root_name+'_pics/'+well_name+'_graph.png', dpi=600, bbox_inches='tight', transparent=True) # hardcode folder and pictures extension # plt.show() 
    # TODO pptx file is too heavy (check if is possible to make lighter pictures) OR SAVE AS SVG (vector format???)
    image_stream = io.BytesIO()
    plt.savefig(image_stream, dpi=100, bbox_inches='tight', format='jpg', optimize=True, quality=70)
    #plt.savefig(image_stream, dpi=100, bbox_inches='tight', format='png')
    make_slide(prs, well_name, well_stat[2], image_stream)

    plt.close()



def make_slide(prs, well_name, well_matched, image_stream):
    slide = prs.slides.add_slide(prs.slide_layouts[TEMPLATE_NUMBER]) # layout number from Slide Master (
    #pic = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(0.0), Inches(1.138), height=Inches(5.098), width=Inches(9.01)) # picture position hardcoded 5.224 height with plt.title
    pic = slide.shapes.add_picture(image_stream, Inches(0.0), Inches(1.138), height=Inches(5.098)) # picture position hardcoded 5.224 height with plt.title

    slide_header_box = slide.shapes.add_textbox(Inches(0.465), Inches(0.441), Inches(8.898), Inches(0.449))
    slide_header = slide_header_box.text_frame
    p = slide_header.paragraphs[0]
    run = p.add_run()
    run.font.name = 'Tahoma'
    run.font.size = Pt(22)
    run.text = well_name

    conclusion_box = slide.shapes.add_textbox(Inches(0.118), Inches(6.468), Inches(9.764), Inches(0.405))
    conclusion = conclusion_box.text_frame
    p = conclusion.paragraphs[0]
    run = p.add_run()
    run.font.name = 'Tahoma'
    run.font.size = Pt(18)
    # TODO need to check well statistics and write correct conclusion
    if well_matched:
        run.text = "Скважина соответствует ТЗ"
        run.font.color.rgb = RGBColor(0, 176, 80)
    else:
        run.text = "Скважина не соответствует ТЗ"
        run.font.color.rgb = RGBColor(228, 108, 10)

    comment_box = slide.shapes.add_textbox(Inches(7.067), Inches(1.291), Inches(2.657), Inches(1.953))
    comment = comment_box.text_frame
    p = comment.paragraphs[0]
    run = p.add_run()
    run.font.name = 'Tahoma'
    run.font.size = Pt(12)
    run.text = "Комментарий"


#TODO find out how to make more beautiful map 
#TODO and insert it into pptx by request
def make_press_map(well_names, well_coordinates, well_pressures):
    class HandlerCircle(HandlerPatch):
        def create_artists(self, legend, orig_handle,
                           xdescent, ydescent, width, height, fontsize, trans):
            r = 5
            x = r + width/2
            y = height/2
            p = mpatches.Circle(xy= (x,y), radius = r)
            self.update_prop(p, orig_handle, legend)
            p.set_transform(trans)
            return [p]
    
    def makeLegendItems(color_list, annotation_list):
        legend_items = []
        for l in range(0, len(color_list)):
            legend_items.append(mpatches.Circle((0, 0),1, fc = color_list[l], edgecolor='black'))
        return (legend_items, annotation_list)

    bubble_scale = 20
    x = [xy[0] for xy in well_coordinates]
    y = [xy[1] for xy in well_coordinates]
    z = [ 0 if p[0]=='-' or p[1]=='-' else float(p[1])-float(p[0]) for p in well_pressures]
    label = [name.strip() for name in well_names]

    fig, (ax1, ax2) = plt.subplots(nrows = 1, ncols = 2, figsize = (10, 7))

    # map1 Plus-Minus map
    plus_minus_colors = ['coral' if d>=0 else 'lightblue' for d in z] 
    sct = ax1.scatter(x, y, s = [abs(i)*bubble_scale for i in z], c = plus_minus_colors, edgecolors = 'black', cmap="Spectral")
    ax1.set_aspect('equal')
    ax1.axes.get_xaxis().set_visible(False)
    ax1.axes.get_yaxis().set_visible(False)
    for lb, xx, yy in zip(label, x, y):
        ax1.annotate(lb, (xx, yy), fontsize = 6)
    legend_items = makeLegendItems(['lightblue', 'coral'], ['< 0', '> 0'])
    ax1.legend( legend_items[0], legend_items[1],  loc = 'upper left', handler_map = {mpatches.Circle: HandlerCircle()})
    #ax1.set_title("Plus-Minus map")


    # map2 Absolute Error map
    error_colors = ['white' if abs(d)< 5 else 'lime' if abs(d)<10  else 'red' for d in z] 
    ax2.scatter(x, y, s = [abs(i)*bubble_scale for i in z], c = error_colors, edgecolors = 'black')
    ax2.set_aspect('equal')
    ax2.axes.get_xaxis().set_visible(False)
    ax2.axes.get_yaxis().set_visible(False)
    for lb, xx, yy in zip(label, x, y):
        ax2.annotate(lb, (xx, yy), fontsize = 6)
    legend_items = makeLegendItems(['white', 'lime', 'red'], ['< 5 bar','5 - 10 bar', '> 10 bar'])
    ax2.legend(legend_items[0], legend_items[1], loc = 'upper left', handler_map = {mpatches.Circle: HandlerCircle()})
    #ax2.set_title("Absolute Error map")

    plt.subplots_adjust(left = 0.01, bottom = 0.01, right = 0.99, top = 0.93, wspace = 0.05)
    plt.savefig('press_bubble_map.jpg', dpi=150, bbox_inches='tight', optimize=True, quality=70)




def get_graphs(currDir, rootName, start_date_array, times, numsArray, RateOut):
    T = len(times)
    MZ = numsArray[5-1]                 # number of connections
    V = cts.VEC + MZ*2*numsArray[55-1]  # number of vectors
    s_d = datetime(start_date_array[2], start_date_array[1], start_date_array[0])

    x_values = [s_d+timedelta(days=x.tos) for x in times]
    y_values = {} 

    #filtered_wells = [e for f in STAT_FILTER for e in RateOut[1]  if re.match(f, e)  and not re.match('WQ2-137',e)]
    #filtered_coords = [e[1] for f in STAT_FILTER for e in list(zip(RateOut[1], RateOut[3])) if re.match(f, e[0]) and not re.match('WQ2-137',e[0])]
    filtered_wells = [e for e in RateOut[1] if re.match(STAT_FILTER, e)]
    filtered_coords = [e[1] for e in list(zip(RateOut[1], RateOut[3])) if re.match(STAT_FILTER, e[0])]

    stat_from_excel = get_statistics(STAT_FILE, [x.strip() for x in filtered_wells])
    stat_table_header = stat_from_excel[0]
    statistics_data = stat_from_excel[1]

    prs = Presentation(PRS_TEMPLATE)
    #Path(rootName+"_pics").mkdir(parents=True, exist_ok=True)

    #filtered_wells = [e for f in STAT_FILTER for e in RateOut[1]  if re.match(f, e)]
    for well_name in filtered_wells:
        wi = RateOut[1].index(well_name) 
        y_values.clear()
        y_values = {'Sbhp':[], 'Hbhp':[], 'Sliq':[], 'Hliq':[], 'Swcut':[], 'Hwcut':[], 'Swir':[], 'Hwir':[], 'wPI4':[]}
        for i in range(T):
            sim_liq =  RateOut[0][T*V*wi+T*cts.i_d['Sopr']+i]+RateOut[0][T*V*wi+T*cts.i_d['Swpr']+i]
            hist_liq = RateOut[0][T*V*wi+T*cts.i_d['Hopr']+i]+RateOut[0][T*V*wi+T*cts.i_d['Hwpr']+i]

            sim_wcut = 0
            if  RateOut[0][T*V*wi+T*cts.i_d['Sopr']+i]==0 and RateOut[0][T*V*wi+T*cts.i_d['Swpr']+i]== 0: sim_wcut = 0
            else: sim_wcut =  RateOut[0][T*V*wi+T*cts.i_d['Swpr']+i]/sim_liq*100

            hist_wcut = 0
            if  RateOut[0][T*V*wi+T*cts.i_d['Hopr']+i]==0 and RateOut[0][T*V*wi+T*cts.i_d['Hwpr']+i]== 0: hist_wcut = 0
            else: hist_wcut =  RateOut[0][T*V*wi+T*cts.i_d['Hwpr']+i]/sim_liq*100

            y_values['Sbhp'].append(RateOut[0][T*V*wi+T*cts.i_d['Sbhp']+i]) # simulated BHP
            y_values['Hbhp'].append(RateOut[0][T*V*wi+T*cts.i_d['Hbhp']+i]) # history BHP
            y_values['Sliq'].append(sim_liq) # simulated Liquid production
            y_values['Hliq'].append(hist_liq) # history Liquid production
            y_values['Swcut'].append(sim_wcut) # simulated WCUT
            y_values['Hwcut'].append(hist_wcut) # history WCUT 
            y_values['Swir'].append(RateOut[0][T*V*wi+T*cts.i_d['Swir']+i]) # simulated Water injection
            y_values['Hwir'].append(RateOut[0][T*V*wi+T*cts.i_d['Hwir']+i]) # history Water injection 
            y_values['wPI4'].append(RateOut[0][T*V*wi+T*cts.i_d['wPI4']+i]) # simulated well production index 4-point 

        make_graph(rootName, well_name.strip(), x_values, y_values, stat_table_header, statistics_data[well_name.strip()], prs)

    make_press_map(filtered_wells, filtered_coords, [statistics_data[e][0][18:20] for f in filtered_wells for e in statistics_data if e == f.strip()])
    prs.save(rootName+'_allWells.pptx')
